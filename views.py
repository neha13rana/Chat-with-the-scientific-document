# chatapp/views.py
from django.shortcuts import render
from django.views.decorators.csrf import csrf_exempt
from django.http import JsonResponse
from django.conf import settings
from django.middleware.csrf import get_token
from django.http import JsonResponse
from django.http import HttpResponse
from PyPDF2 import PdfReader
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import OpenAIEmbeddings
from langchain_community.llms import OpenAI
from langchain.chains.question_answering import load_qa_chain
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
# from langchain.vectorstores import FAISS
# from langchain.embeddings.openai import OpenAIEmbeddings
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain.chains.question_answering import load_qa_chain
# from langchain.llms import OpenAI
# from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document

def chat_view(request):
    return render(request, 'chatapp/chat.html')

def read_pdf_file(file_path):
    with open(file_path, 'rb') as file:
        pdf_reader = PdfReader(file)
        text = ''
        for page in pdf_reader.pages:
            text += page.extract_text() + '\n'
        return text.strip()
def extract_text_from_pptx(pptx_path):
    text = ""
    presentation = Presentation(pptx_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
def extract_text_from_ppt(ppt_path):
    text = ""
    presentation = Presentation(ppt_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text_frame.text + "\n"
    return text
def extract_text_from_docx(file_path):
    doc = Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    text =  '\n'.join(full_text)
    return text
def get_file_type(file_name):
    _, file_extension = os.path.splitext(file_name)
    return file_extension.lower()
def extract_text_from_file(file_name):
    file_type = get_file_type(file_name)
    if file_type == '.pdf':
        return read_pdf_file(file_name)
    elif file_type == '.pptx':
        return extract_text_from_pptx(file_name)
    elif file_type == '.ppt':
        return extract_text_from_ppt(file_name)
    elif file_type == '.docx':
        return extract_text_from_docx(file_name)
    else:
        return "Unsupported file type."

def text_to_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(
              chunk_size=512,
              chunk_overlap=32,
              length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks
def processing_file(chunks,file_name):
  os.environ["OPENAI_API_KEY"] = "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"
  embeddings = OpenAIEmbeddings(openai_api_key="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")
  docsearch = FAISS.from_texts(chunks, embeddings)
  chain = load_qa_chain(OpenAI(), chain_type="stuff")
  print(f"\nProcessing PDF file: {file_name}\n")
  return docsearch, chain
def index(request):
    if request.method == 'POST' and request.FILES['file']:
        uploaded_file = request.FILES['file']
        text = extract_text_from_file(uploaded_file)
        chunks = text_to_chunks(text)
        docsearch, chain = processing_file(chunks, uploaded_file.name)
        return render(request, 'result.html', {'docsearch': docsearch, 'chain': chain})
    else:
        form = UploadFileForm()
        return render(request, 'upload.html', {'form': form})

def question_answering(request):
    if request.method == 'POST':
        question = request.POST['question']
        docs = request.POST['docs']
        chain = request.POST['chain']
        
        if chain and hasattr(chain, 'run'):
            answer = chain.run(input_documents=docs, question=question)
            return HttpResponse(answer)
        else:
            return HttpResponse("Chain object not found or does not have the run method.")
    return HttpResponse("Invalid request")

def upload_file(request):
    if request.method == 'POST' and request.FILES['file']:
        uploaded_file = request.FILES['file']
        file_path = os.path.join(settings.MEDIA_ROOT, uploaded_file.name)
        with open(file_path, 'wb+') as destination:
            for chunk in uploaded_file.chunks():
                destination.write(chunk)

        text = extract_text_from_file(file_path)
        chunks = text_to_chunks(text)
        docsearch, chain = processing_file(chunks, uploaded_file.name)
        question_answering(docsearch, chain)

        # You can customize this response as needed
        return HttpResponse("File uploaded and processed successfully.")

    return render(request, 'upload.html', {})
@csrf_exempt
def process_document(request):
    if request.method == 'POST':
        # Get the uploaded PDF file from the request
        uploaded_file = request.FILES.get('document')
        if not uploaded_file:
            return JsonResponse({'error': 'No document file uploaded'}, status=400)

        return JsonResponse({'answers': answers})
    else:
        return JsonResponse({'error': 'Invalid request method'}, status=405)
